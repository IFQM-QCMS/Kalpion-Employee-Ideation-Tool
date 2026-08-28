# -*- coding: utf-8 -*-
"""
Build IFQM_Employee_Ideation_Tool_Overview.pptx from V2_DWM_Overview.pptx.

    python "Customer Presentation/build_deck.py"

── Why this copies the file instead of drawing a new one ────────────────────

The source deck is hand-built: every panel, chevron and card is an individually
positioned autoshape, most of them carrying gradient fills that no reasonable
amount of code would reproduce faithfully. Rebuilding it would produce
something that resembled the design rather than something that IS it.

So the file is copied and rewritten in place. The theme, the slide master, both
layouts, the title slide's artwork, the purple spine, the logo, every gradient
and every rounded corner survive untouched, because nothing here creates or
deletes a shape — it only changes the words inside them and the bytes behind
the pictures.

That is possible because the two decks turn out to have the same shape: an
opening, a "why", a "what is it", a role/permission matrix, an architecture
slide, and four screen tours. Each slide of the original maps onto one of ours.

── How text is replaced without losing formatting ───────────────────────────

Setting `shape.text_frame.text` would drop every run and take its font, size
and colour with it. Instead `retext()` writes into the FIRST run and empties the
rest, so the run properties the designer set — Poppins Medium 24pt purple for a
heading, Calibri 12.5pt orange for a kicker — carry over to the new words.

── The two structural edits ─────────────────────────────────────────────────

Slide 9 of the source ends with a full-bleed picture sitting on top of its own
content, hiding the four cards underneath it. That overlay is removed so the
designed slide is visible. Slide 4's page number reads "5" and slide 9's reads
"7" in the original; the numbers are renumbered 1-9 to match their new order.
"""
import copy
import os
import shutil
import zipfile

from pptx import Presentation

HERE = os.path.dirname(os.path.abspath(__file__))
SRC = os.path.join(HERE, 'V2_DWM_Overview.pptx')
OUT = os.path.join(HERE, 'IFQM_Employee_Ideation_Tool_Overview.pptx')
ART = os.path.join(HERE, 'art')


# ── Text helpers ────────────────────────────────────────────────────────────

def retext(shape, value):
    """
    Replace a shape's words, keeping the designer's run formatting.

    A plain string goes into the first run; the rest are emptied rather than
    deleted, so paragraph spacing is undisturbed. Setting text_frame.text
    instead would drop every run and take its font, size and colour with it.

    A TUPLE writes one entry per run, which several shapes need because their
    formatting lives in the run split rather than the paragraph. The value box
    on slide 3 is "~15%" at 30pt bold white followed by "operational
    efficiency" at 16pt pale blue — collapsing that into one run renders the
    whole line at 30pt, which wraps to three lines and climbs over its own
    heading. The banner and the two lead-in paragraphs on slide 2 are the same
    shape of thing: a bold opening clause, then a lighter continuation.

    Extra runs beyond the tuple are emptied; a tuple longer than the available
    runs raises, rather than silently dropping the tail.
    """
    tf = shape.text_frame
    paras = [p for p in tf.paragraphs]
    if not paras:
        return
    first = paras[0]
    if not first.runs:
        first.add_run()

    if isinstance(value, tuple):
        runs = [r for r in first.runs]
        # Runs holding only whitespace are separators the designer put between
        # the styled pieces; they are kept as-is and not counted.
        styled = [r for r in runs if r.text.strip()]
        if len(value) > len(styled):
            raise ValueError('%s has %d styled run(s), got %d values'
                             % (shape.name, len(styled), len(value)))
        for r, text in zip(styled, value):
            r.text = text
        for r in styled[len(value):]:
            r.text = ''
    else:
        first.runs[0].text = value
        for r in first.runs[1:]:
            r.text = ''

    for p in paras[1:]:
        for r in p.runs:
            r.text = ''


def by_name(slide):
    return {sh.name: sh for sh in slide.shapes}


def apply(slide, mapping):
    """Rewrite every named shape on this slide. Unknown names are reported."""
    shapes = by_name(slide)
    missing = [k for k in mapping if k not in shapes]
    if missing:
        raise KeyError('shapes not on slide: %s' % missing)
    for name, value in mapping.items():
        retext(shapes[name], value)


def drop(slide, name):
    sh = by_name(slide).get(name)
    if sh is not None:
        sh._element.getparent().remove(sh._element)


# ── Slide content ───────────────────────────────────────────────────────────

S1 = {
    'Title 5': 'Kalpion',
    'Subtitle 6': 'A digital platform that turns everyday employee ideas into measurable improvement.',
}

S2 = {
    'Text 0': 'Why Kalpion?',
    'Text 6': 'WHY WE ARE DOING THIS',
    'Text 7': ('The people closest to the work see the waste first',
               ' — a jig that slips, a rework loop, a delivery that always runs late. Most of '
               'it is never written down, and what is written down goes into a suggestion box '
               'nobody empties.'),
    'Text 9': ('The Result: ',
               'Every idea captured in one place, routed to the person who can decide, and '
               'tracked through to what it actually saved.'),
    'Text 10': ('Objective — ',
                'Give every employee a way to raise an improvement, every manager a queue they '
                'can act on, and every plant a record of what was implemented and what it was '
                'worth.'),
    'Text 11': 'WHAT CHANGES WITH THE PLATFORM',
    'Text 12': 'CHALLENGE TODAY',
    'Text 13': 'WITH THE PLATFORM',
    'Text 16': 'Ideas lost in a suggestion box',
    'Text 20': 'One place, tracked end to end',
    'Text 23': 'Nobody knows who should approve',
    'Text 27': 'A named, ordered approval chain',
    'Text 30': 'The employee never hears back',
    'Text 34': 'Status, comments and points',
    'Text 37': 'Savings are never measured',
    'Text 41': 'ROI recorded on every idea',
}

S3 = {
    'Text 0': 'What is Kalpion?',
    'Text 7': ('A platform where employees raise improvement ideas and somebody actually '
               'decides on them.',
               '  Ideas move up a named approval chain, and approved ones are pushed to the QC '
               'tool and tracked to implementation — with the benefit recorded.'),
    'Text 10': 'Capture',
    'Text 11': ('An employee describes the problem, the proposed solution and the expected '
                'benefit, and attaches evidence.'),
    'Text 14': 'Evaluate',
    'Text 15': ('The idea is scored, discussed and moved up the chain until somebody with the '
                'authority decides.'),
    'Text 18': 'Implement',
    'Text 19': ('Approved ideas go to the QC tool as tracked work, and their ROI is recorded '
                'against the original idea.'),
    'Text 21': 'How it works',
    'Text 22': 'Employee submits an idea',
    'Text 24': 'Scored and reviewed',
    'Text 26': 'Approved up the chain',
    'Text 28': 'Implemented and measured',
    'Text 30': 'WHAT YOU GET OUT OF THE BOX',
    'Text 31': ('7', 'languages, every screen'),
    'Text 32': ('6', 'stages, Draft to Implemented'),
}

S4 = {
    'Text 0': 'Role-Wise Flow, Data Entry and Reports',
    'Text 1': 'Four user levels, what each one does, and how information moves between them',
    'Text 3': 'Role',
    'Text 5': 'What they do',
    'Text 7': 'Data entry',
    'Text 9': 'Reports',
    'Text 26': 'Setup flows down',
    'Text 27': 'Ideas flow up',
    'Text 28': 'tenant isolation',

    # Level 1 — platform
    'Text 30': 'IFQM Platform Admin',
    'Text 31': 'Level 1 · Platform',
    'Text 34': 'Onboard Org',
    'Text 37': 'Assign Plan',
    'Text 40': 'Set Policy',
    'Text 42': 'Organisations · Plans',
    'Text 46': 'Usage and billing only',
    'Text 48': 'Never reads an idea',

    # Level 2 — organisation
    'Text 50': 'Organisation Admin',
    'Text 51': 'Level 2 · Organisation',
    'Text 54': 'Add Users',
    'Text 57': 'Set Approval Chain',
    'Text 60': 'Configure Branding',
    'Text 64': 'Users · Roles · Departments',
    'Text 66': 'Approval chain · Settings',
    'Text 68': 'Org-wide analytics',
    'Text 70': 'PDF export',

    # Level 3 — reviewers
    'Text 72': 'Manager / Plant Head',
    'Text 73': 'Level 3 · Approval chain',
    'Text 76': 'Open Review Queue',
    'Text 79': 'Score & Comment',
    'Text 82': 'Approve or Escalate',
    'Text 86': 'Decisions · Comments',
    'Text 88': 'ROI · Implementation',
    'Text 90': 'Team and department view',
    'Text 92': 'PDF export',

    # Level 4 — employees
    'Text 94': 'Employee',
    'Text 95': 'Level 4 · Shop floor',
    'Text 98': 'Sign In (OTP)',
    'Text 101': 'Submit Idea',
    'Text 104': 'Track Status',
    'Text 108': 'Idea · Benefit · Files',
    'Text 110': 'Comments · Votes',
    'Text 112': 'Own ideas and points',
    'Text 114': 'No export',

    'Text 115': ('Setup and configuration flow downward · ideas and their outcomes flow upward, '
                 'stopping at the boundary of each organisation’s own database.'),
}

S5 = {
    'Text 0': 'High-Level System Architecture',
    'Text 6': 'Frontend',
    'Text 7': 'React (Vite) — responsive, role-based UI',
    'Text 9': 'Backend / API',
    'Text 10': 'Node.js + Express REST server',
    'Text 12': 'Database',
    'Text 13': 'MySQL — one database per organisation',
    'Text 15': 'Hosting',
    'Text 16': 'Vercel + Render · Aiven MySQL',
    'Text 18': 'Security',
    'Text 19': 'HTTPS · JWT · bcrypt · tenant isolation',
}

S6 = {
    'Text 0': 'The Employee Experience',
    'Text 7': 'AT A GLANCE',
    'Text 1': ('Everything an employee needs on one screen: raise an idea, follow it up the '
               'chain, and see what it earned.'),
    'Text 10': 'Submit in minutes',
    'Text 13': 'Follow every step',
    'Text 16': 'Points and badges',
    'Text 19': 'Community board',
}

S7 = {
    'Text 0': 'Review and Approval',
    'Text 7': 'WHAT A REVIEWER DOES',
    'Text 1': ('Reviewers see exactly what is waiting on them, how long it has waited, and who '
               'it goes to next.'),
    'Text 10': 'One clear queue',
    'Text 13': 'Score and comment',
    'Text 16': 'Approve or escalate',
    'Text 19': 'Overdue flagged',
}

S8 = {
    'Text 0': 'Organisation Analytics',
    'Text 7': 'WHAT AN ADMIN CAN SEE',
    'Text 1': 'Where ideas come from, how fast they move, and what they are worth.',
    'Text 10': 'Approval rate',
    'Text 13': 'Implementation rate',
    'Text 16': 'Quality score',
    'Text 19': 'Department view',
}

S9 = {
    'Text 0': 'Platform Console — Multi-Tenant',
    'Text 1': ('One console to onboard organisations, set their plan and watch the health of '
               'the platform — without ever reading a customer’s ideas.'),
    'Text 7': 'Platform Admin → Organisations · plans · usage · Export to PDF',
    'Text 8': 'AT A GLANCE',
    'Text 11': 'Organisations',
    'Text 12': 'Onboard, approve, suspend',
    'Text 15': 'Plans & Billing',
    'Text 16': 'Trial · Starter · Professional · Lifetime',
    'Text 19': 'Usage & Quotas',
    'Text 20': 'Requests, users and storage per plan',
    'Text 23': 'Support & Health',
    'Text 24': 'Tickets, maintenance mode, DB checks',
}

SLIDES = [S1, S2, S3, S4, S5, S6, S7, S8, S9]

# Slide index (0-based) -> the media file its content picture uses.
PICTURES = {
    4: 'image9.png',    # architecture
    5: 'image10.png',   # employee dashboard
    6: 'image11.png',   # review queue
    7: 'image12.png',   # analytics
    8: 'image13.png',   # platform console
}


def main():
    shutil.copyfile(SRC, OUT)
    prs = Presentation(OUT)

    for i, mapping in enumerate(SLIDES):
        slide = prs.slides[i]
        apply(slide, mapping)

    # Slide 9 carries a full-bleed picture ON TOP of its own content, hiding the
    # four cards beneath it. Removed so the designed slide is the one that shows.
    drop(prs.slides[8], 'Picture 30')

    # The source numbers its footers 1,2,3,5,6,7,8,9,7 — two are wrong and one
    # repeats. Renumbered to match the order the slides are actually in.
    NUM_SHAPE = {1: 'Text 1', 2: 'Text 1', 3: 'Text 2', 4: 'Text 1',
                 5: 'Text 2', 6: 'Text 2', 7: 'Text 2', 8: 'Text 2'}
    for idx, shape_name in NUM_SHAPE.items():
        sh = by_name(prs.slides[idx]).get(shape_name)
        if sh is not None and sh.has_text_frame:
            retext(sh, str(idx + 1))

    # The copy inherits the source deck's document properties, so a file named
    # for this product would still announce itself as the DWM overview in the
    # title bar, in search results and in any library it is uploaded to.
    cp = prs.core_properties
    cp.title = 'Kalpion — Overview'
    cp.subject = 'Customer presentation'
    cp.comments = ('Built from V2_DWM_Overview.pptx by build_deck.py — same theme, master, '
                   'layouts and design; content and artwork replaced.')

    prs.save(OUT)

    # ── Swap the artwork ──
    #
    # Done on the saved package rather than through python-pptx: rewriting the
    # zip entry keeps every picture frame's position, crop and z-order exactly
    # as the designer left them, because the frame is never touched at all.
    tmp = OUT + '.tmp'
    replacements = {}
    for name in set(PICTURES.values()):
        path = os.path.join(ART, name)
        if not os.path.exists(path):
            raise SystemExit('missing artwork: %s — run build_images.py first' % path)
        with open(path, 'rb') as fh:
            replacements['ppt/media/' + name] = fh.read()

    with zipfile.ZipFile(OUT, 'r') as zin, \
            zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = replacements.get(item.filename, zin.read(item.filename))
            zout.writestr(item, data)
    os.replace(tmp, OUT)

    print('saved', OUT)
    print('slides:', len(prs.slides.__iter__.__self__._sldIdLst))
    for name in sorted(replacements):
        print('  artwork replaced:', name)


if __name__ == '__main__':
    main()
